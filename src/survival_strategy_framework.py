"""
====================================================================================================
TITLE:              SURVIVAL STRATEGY FRAMEWORK
SUBTITLE:           From Time-to-Event Data to Governed Intervention Evidence
AUTHOR:             Andrew R. Goad
PUBLIC REPOSITORY:  github.com/andrew-goad/survival-strategy-framework
ENGINE RELEASE:     2026.07-mainline
PYTHON:             3.11+
----------------------------------------------------------------------------------------------------
FUNCTIONAL PURPOSE
----------------------------------------------------------------------------------------------------
This module implements a reusable, portfolio-grade survival-analysis decision framework. It is
intentionally broader than a one-time churn model or notebook. The framework connects:

    governed time-to-event inputs
        -> independent persona discovery
        -> regularized Cox proportional hazards modeling
        -> model-derived risk stratification
        -> same-cohort intervention simulation
        -> dependency-safe feature reconstruction
        -> predicted survival comparison
        -> archived executive and technical evidence

The primary demonstration uses deterministic synthetic retention data. A secondary Rossi-dataset
smoke test is retained only to prove that the same framework can operate against another conforming
survival dataset.

----------------------------------------------------------------------------------------------------
EXECUTIVE TALK TRACKS
----------------------------------------------------------------------------------------------------
1. PREDICTING THE RUNWAY
   Binary classification asks whether an event may occur. Survival analysis adds the timing
   dimension: when is risk likely to emerge, and which populations have the shortest modeled
   runway?

2. PERSONAS AND RISK ARE COMPLEMENTARY — NOT THE SAME THING
   K-Means personas describe naturally occurring behavioral groups. The CoxPH model independently
   estimates relative event hazard. Persona labels are not inserted into the CoxPH feature matrix.
   This separation prevents the clustering layer from being misrepresented as a model driver.

3. TARGETED, SAME-COHORT SIMULATION
   The framework identifies a governed high-risk cohort from baseline partial-hazard scores, applies
   controlled scenario changes only to that cohort, rebuilds dependent engineered terms, and
   re-scores the identical IDs. This isolates modeled scenario movement from population drift.

4. PUSHING THE CURVE TO THE RIGHT
   A configured scenario is evaluated through both relative partial-hazard movement and modeled
   survival probability at a defined horizon. A favorable scenario should reduce modeled hazard and
   raise the same cohort's predicted survival curve.

5. NO COLD HANDOFFS
   The output is not only a fitted model. Each run creates a reproducible evidence package with input
   validation, model coefficients and hazard ratios, proportional-hazards tests, cohort profiles,
   scenario definitions, dependency audits, survival visuals, an executive PowerPoint, and a
   technical PDF.

----------------------------------------------------------------------------------------------------
STATISTICAL AND GOVERNANCE BOUNDARIES
----------------------------------------------------------------------------------------------------
- This is a survival-analysis and scenario-sensitivity framework, not a causal-inference engine.
- Scenario results describe modeled sensitivity under controlled feature changes. They do not prove
  that an operational intervention will cause the modeled outcome.
- The primary demonstration uses synthetic data and does not contain PII or proprietary customer
  information.
- Partial hazard is a relative-risk quantity. It is not a probability.
- Kaplan-Meier curves are non-parametric observed survival estimates. The baseline-versus-scenario
  curve is a CoxPH model-predicted survival comparison. They are deliberately labeled separately.
- The framework does not calculate booked revenue, realized ROI, or defended LTV.
- Production use requires domain-specific event definitions, censoring rules, feature governance,
  out-of-sample validation, calibration review, stability monitoring, and applicable model-risk
  approval.

----------------------------------------------------------------------------------------------------
CORE DATA CONTRACT
----------------------------------------------------------------------------------------------------
Every input dataset must provide:

1. A unique ID column.
2. A strictly positive numeric duration column.
3. A binary event column where 1 = event observed and 0 = right-censored / still active.
4. Numeric or Boolean model features with no missing or infinite values.
5. Numeric or Boolean segmentation features with no missing or infinite values.
6. Engineered terms that follow these conventions when used:
      - interaction: <feature_a>_x_<feature_b>[_x_<feature_c>...]
      - square:      <feature>_sq

Scenarios may change governed base features. They may not directly change the ID, duration, event,
or engineered-term columns. Engineered terms are rebuilt automatically after base-feature changes.

----------------------------------------------------------------------------------------------------
RUN OUTPUT PACKAGE
----------------------------------------------------------------------------------------------------
Each successful run creates an immutable timestamped directory containing:

    run_metadata.json
    framework_config.json
    scenario_definitions.json
    input_validation.csv
    input_snapshot.csv                         (optional; enabled in the public demo)
    model_summary.csv
    proportional_hazards_test.csv
    model_fit_warnings.csv
    model_scaler_parameters.csv
    cox_baseline_survival.csv
    cox_baseline_cumulative_hazard.csv
    persona_scaler_parameters.csv
    persona_centroids_standardized.csv
    persona_centroids_raw_scale.csv
    persona_profiles.csv
    risk_tier_profiles.csv
    target_cohort.csv
    scenario_results.csv
    scenario_target_scores.csv
    dependency_audit.csv
    predicted_survival_curves.csv
    acceptance_checks.csv
    executive_narrative.txt
    charts/
        persona_kaplan_meier.png
        risk_tier_kaplan_meier.png
        baseline_vs_scenario_survival.png
        scenario_hazard_reduction.png
    Survival_Strategy_Deck.pptx
    Technical_Model_Evidence.pdf
    artifact_manifest.json

----------------------------------------------------------------------------------------------------
PRIMARY DEPENDENCIES
----------------------------------------------------------------------------------------------------
- pandas / numpy: data engineering and controlled simulation
- scikit-learn: StandardScaler and reproducible K-Means personas
- lifelines: CoxPH, Kaplan-Meier, predicted survival, and proportional-hazards testing
- matplotlib: evidence visuals
- python-pptx: executive reporting
- reportlab: technical model-evidence PDF
====================================================================================================
"""

from __future__ import annotations

# ==================================================================================================
# STANDARD LIBRARY IMPORTS
# ==================================================================================================

import argparse
import hashlib
import importlib.metadata
import json
import platform
import sys
import traceback
import warnings
from dataclasses import asdict, dataclass, field, is_dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Literal, Mapping, Sequence

# ==================================================================================================
# THIRD-PARTY ANALYTICAL IMPORTS
# ==================================================================================================

import matplotlib

# Headless rendering is required for servers, GitHub Actions, and command-line execution.
matplotlib.use("Agg")

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from sklearn.cluster import KMeans
from sklearn.preprocessing import StandardScaler

# Survival-analysis engine.
from lifelines import CoxPHFitter, KaplanMeierFitter
from lifelines.datasets import load_rossi
from lifelines.statistics import proportional_hazard_test

# Executive-reporting stack.
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Technical-reporting stack.
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import landscape, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image as RLImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


# ==================================================================================================
# ENGINE CONSTANTS AND VISUAL TOKENS
# ==================================================================================================

ENGINE_RELEASE = "2026.07-mainline"

# Report colors intentionally mirror the broader portfolio's restrained dark technical aesthetic.
COLOR_BACKGROUND = "#07131A"
COLOR_PANEL = "#0C202A"
COLOR_CYAN = "#20D7E5"
COLOR_TEAL = "#1BA7A8"
COLOR_GOLD = "#D9A441"
COLOR_WHITE = "#F4F5F6"
COLOR_LIGHT_GRAY = "#C7D0D5"
COLOR_MID_GRAY = "#71808A"
COLOR_RED = "#D95C5C"
COLOR_GREEN = "#56C271"


# ==================================================================================================
# CONFIGURATION CONTRACTS
# ==================================================================================================

@dataclass(frozen=True)
class FrameworkConfig:
    """Defines the governed analytical and reporting contract for one framework run.

    Why a configuration object matters:
    - It separates assumptions from modeling logic.
    - It makes runs reproducible and reviewable.
    - It prevents scenario definitions and model controls from being hidden inside functions.
    - It allows the same orchestrator to support the synthetic retention demo and Rossi smoke test.
    """

    project_name: str
    duration_col: str
    event_col: str
    id_col: str

    # Persona and CoxPH features are deliberately separate. Persona labels do not enter the model.
    model_features: tuple[str, ...]
    segmentation_features: tuple[str, ...]

    # Optional profile fields displayed in persona / risk-tier summaries.
    profile_features: tuple[str, ...] = ()

    n_clusters: int = 3
    risk_quantile: float = 0.75
    penalizer: float = 0.10
    l1_ratio: float = 0.0
    random_state: int = 42

    # CoxPH features are standardized using a scaler fitted only on the baseline population.
    # Scenario data use that same fitted scaler so baseline and scenario scores remain comparable.
    standardize_model_features: bool = True

    # Technical rigor controls.
    show_km_confidence_intervals: bool = True
    run_ph_test: bool = True
    ph_test_alpha: float = 0.05
    evaluation_horizon: float = 12.0
    timeline_points: int = 121

    # Validation thresholds are warnings unless the underlying data violate a hard contract.
    low_event_rate_warning: float = 0.05
    high_event_rate_warning: float = 0.95
    high_missingness_warning: float = 0.20

    # Output governance.
    output_root: Path = Path("outputs")
    save_input_snapshot: bool = True


@dataclass(frozen=True)
class FeatureChange:
    """Defines one governed scenario change to a base feature.

    operation:
        multiply -> new = old * value
        add      -> new = old + value
        replace  -> new = value

    lower_bound / upper_bound protect the domain of the changed feature.
    round_digits is useful for count-like fields such as product adoption.
    """

    feature: str
    operation: Literal["multiply", "add", "replace"]
    value: float
    lower_bound: float | None = None
    upper_bound: float | None = None
    round_digits: int | None = None


@dataclass(frozen=True)
class ScenarioDefinition:
    """Defines one scenario and the assumptions that distinguish it from baseline."""

    name: str
    description: str
    changes: tuple[FeatureChange, ...] = field(default_factory=tuple)


# ==================================================================================================
# RESULT AND ARTIFACT CONTAINERS
# ==================================================================================================

@dataclass(frozen=True)
class ValidationIssue:
    severity: Literal["ERROR", "WARNING", "INFO"]
    area: str
    field: str
    message: str


@dataclass
class ValidationResult:
    issues: list[ValidationIssue] = field(default_factory=list)

    def add(self, severity: str, area: str, field_name: str, message: str) -> None:
        self.issues.append(
            ValidationIssue(
                severity=severity.upper(),  # type: ignore[arg-type]
                area=area,
                field=field_name,
                message=message,
            )
        )

    @property
    def has_errors(self) -> bool:
        return any(issue.severity == "ERROR" for issue in self.issues)

    def to_frame(self) -> pd.DataFrame:
        if not self.issues:
            return pd.DataFrame(
                [{
                    "Severity": "INFO",
                    "Area": "Input Validation",
                    "Field": "ALL",
                    "Message": "All hard input-contract checks passed.",
                }]
            )
        return pd.DataFrame(
            [
                {
                    "Severity": issue.severity,
                    "Area": issue.area,
                    "Field": issue.field,
                    "Message": issue.message,
                }
                for issue in self.issues
            ]
        )

    def raise_for_errors(self) -> None:
        if not self.has_errors:
            return
        error_text = "\n".join(
            f"[{issue.area}] {issue.field}: {issue.message}"
            for issue in self.issues
            if issue.severity == "ERROR"
        )
        raise ValueError(f"Survival input validation failed:\n{error_text}")


@dataclass
class PersonaArtifacts:
    scored_data: pd.DataFrame
    scaler: StandardScaler
    kmeans: KMeans
    profiles: pd.DataFrame


@dataclass
class CoxModelArtifacts:
    model: CoxPHFitter
    model_scaler: StandardScaler | None
    training_frame: pd.DataFrame
    model_summary: pd.DataFrame
    fit_metadata: pd.DataFrame
    ph_test_results: pd.DataFrame
    fit_warnings: pd.DataFrame


@dataclass
class RiskArtifacts:
    scored_data: pd.DataFrame
    risk_tier_profiles: pd.DataFrame
    target_cohort: pd.DataFrame
    target_index: pd.Index
    target_threshold: float
    target_label: str


@dataclass
class ScenarioArtifacts:
    summary: pd.DataFrame
    target_scores: pd.DataFrame
    dependency_audit: pd.DataFrame
    predicted_survival_curves: pd.DataFrame


@dataclass
class PlotArtifacts:
    persona_km: Path
    risk_tier_km: Path
    baseline_vs_scenario: Path
    scenario_comparison: Path


@dataclass
class RunArtifacts:
    run_id: str
    run_directory: Path
    validation: pd.DataFrame
    personas: PersonaArtifacts
    model: CoxModelArtifacts
    risk: RiskArtifacts
    scenarios: ScenarioArtifacts
    plots: PlotArtifacts
    executive_narrative: str
    acceptance_checks: pd.DataFrame
    power_point: Path
    technical_pdf: Path
    manifest: Mapping[str, str]


# ==================================================================================================
# GENERIC UTILITY FUNCTIONS
# ==================================================================================================

def utc_now() -> datetime:
    """Returns a timezone-aware UTC timestamp for reproducible run lineage."""
    return datetime.now(timezone.utc)


def safe_slug(value: str) -> str:
    """Converts free text into a file-system-safe identifier."""
    cleaned = "".join(ch.lower() if ch.isalnum() else "-" for ch in value.strip())
    while "--" in cleaned:
        cleaned = cleaned.replace("--", "-")
    return cleaned.strip("-") or "run"


def to_jsonable(value: Any) -> Any:
    """Recursively converts dataclasses, Paths, tuples, and NumPy values into JSON-safe values."""
    if is_dataclass(value):
        return {key: to_jsonable(item) for key, item in asdict(value).items()}
    if isinstance(value, Path):
        return str(value)
    if isinstance(value, Mapping):
        return {str(key): to_jsonable(item) for key, item in value.items()}
    if isinstance(value, (tuple, list)):
        return [to_jsonable(item) for item in value]
    if isinstance(value, (np.integer,)):
        return int(value)
    if isinstance(value, (np.floating,)):
        return float(value)
    if isinstance(value, (np.bool_,)):
        return bool(value)
    if isinstance(value, datetime):
        return value.isoformat()
    return value


def write_json(path: Path, payload: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as stream:
        json.dump(to_jsonable(payload), stream, indent=2, sort_keys=True)


def package_versions() -> dict[str, str]:
    """Captures the execution environment without failing if a package name is unavailable."""
    packages = [
        "pandas",
        "numpy",
        "matplotlib",
        "scikit-learn",
        "lifelines",
        "python-pptx",
        "reportlab",
    ]
    versions: dict[str, str] = {"python": platform.python_version()}
    for package in packages:
        try:
            versions[package] = importlib.metadata.version(package)
        except importlib.metadata.PackageNotFoundError:
            versions[package] = "NOT_INSTALLED"
    return versions


def dataframe_fingerprint(df: pd.DataFrame, id_col: str) -> str:
    """Creates a stable SHA-256 fingerprint for the run input.

    Sorting by the governed ID prevents a harmless row-order change from producing a different
    fingerprint. The fingerprint is lineage evidence; it is not a cryptographic privacy control.
    """
    ordered = df.sort_values(id_col).reset_index(drop=True)
    hashed = pd.util.hash_pandas_object(ordered, index=False).values.tobytes()
    return hashlib.sha256(hashed).hexdigest()


def create_run_directory(config: FrameworkConfig, fingerprint: str) -> tuple[str, Path]:
    timestamp = utc_now().strftime("%Y%m%dT%H%M%SZ")
    run_id = f"{timestamp}_{fingerprint[:8]}"
    run_directory = config.output_root / safe_slug(config.project_name) / run_id
    (run_directory / "charts").mkdir(parents=True, exist_ok=False)
    return run_id, run_directory


def numeric_frame(df: pd.DataFrame, columns: Sequence[str]) -> pd.DataFrame:
    """Returns a clean floating-point feature matrix while preserving row index and column names."""
    result = df.loc[:, list(columns)].copy()
    for column in result.columns:
        if pd.api.types.is_bool_dtype(result[column]):
            result[column] = result[column].astype(int)
        result[column] = pd.to_numeric(result[column], errors="raise").astype(float)
    return result


def zscore(values: np.ndarray) -> np.ndarray:
    """Stable z-score helper used only inside deterministic synthetic-data construction."""
    std = float(np.std(values))
    if std == 0:
        return np.zeros_like(values, dtype=float)
    return (values - float(np.mean(values))) / std


def format_percent(value: float, decimals: int = 2) -> str:
    if pd.isna(value):
        return ""
    return f"{value:.{decimals}%}"


def format_float(value: float, decimals: int = 3) -> str:
    if pd.isna(value):
        return ""
    return f"{value:,.{decimals}f}"


# ==================================================================================================
# DETERMINISTIC SYNTHETIC RETENTION DEMONSTRATION
# ==================================================================================================

def generate_synthetic_retention_data(
    n_records: int = 7_500,
    seed: int = 42,
    max_followup_months: float = 24.0,
) -> pd.DataFrame:
    """Generates deterministic synthetic customer-lifecycle data.

    Design principles:
    - No real customer data or PII.
    - Independent censoring is generated separately from event timing.
    - Event timing follows a proportional-hazards-style exponential process whose relative hazard
      depends on engagement, adoption, service burden, tenure, and renewal characteristics.
    - Relationships are directional, not deterministic. Mixed-signal records remain present.
    - Engineered interaction and squared terms are included so dependency-safe simulation can be
      demonstrated and independently audited.

    The resulting dataset is a modeling demonstration. Its distributions and coefficients are not
    institution-calibrated retention estimates.
    """
    if n_records < 500:
        raise ValueError("n_records must be at least 500 for stable segmentation and survival modeling.")
    if max_followup_months <= 1:
        raise ValueError("max_followup_months must be greater than 1.")

    rng = np.random.default_rng(seed)

    customer_id = np.array([f"CUST-{i:07d}" for i in range(1, n_records + 1)])
    enterprise_flag = rng.binomial(1, 0.28, n_records)

    # Latent relationship health introduces realistic correlation without becoming a model feature.
    latent_relationship_health = rng.normal(0.0, 1.0, n_records)

    onboarding_completion = np.clip(
        rng.beta(4.5, 2.2, n_records)
        + 0.05 * enterprise_flag
        + 0.04 * latent_relationship_health,
        0.05,
        1.00,
    )

    monthly_active_days = np.clip(
        rng.normal(
            8.0
            + 13.0 * onboarding_completion
            + 1.5 * enterprise_flag
            + 1.2 * latent_relationship_health,
            3.5,
            n_records,
        ),
        0.0,
        30.0,
    )

    product_adoption_count = np.clip(
        np.rint(
            rng.poisson(
                0.7
                + 2.3 * onboarding_completion
                + 0.08 * monthly_active_days
                + 0.4 * enterprise_flag
            )
        ),
        1,
        8,
    ).astype(int)

    support_tickets_90d = np.clip(
        rng.poisson(
            0.8
            + 1.8 * (1.0 - onboarding_completion)
            + 0.35 * np.maximum(0.0, -latent_relationship_health)
        ),
        0,
        12,
    ).astype(int)

    support_resolution_days = np.clip(
        rng.gamma(shape=2.0, scale=1.2, size=n_records)
        + 0.45 * support_tickets_90d
        - 0.8 * enterprise_flag,
        0.25,
        20.0,
    )

    service_incidents_90d = np.clip(
        rng.poisson(
            0.25
            + 0.18 * support_tickets_90d
            + 0.04 * support_resolution_days
        ),
        0,
        8,
    ).astype(int)

    tenure_at_start = np.clip(
        rng.gamma(shape=2.2, scale=5.5, size=n_records) + 3.0 * enterprise_flag,
        0.5,
        60.0,
    )

    renewal_logit = (
        -0.4
        + 1.6 * onboarding_completion
        + 0.08 * monthly_active_days
        + 0.25 * product_adoption_count
        - 0.08 * support_resolution_days
        + 0.4 * enterprise_flag
    )
    auto_renew_probability = 1.0 / (1.0 + np.exp(-renewal_logit))
    auto_renew_flag = rng.binomial(1, np.clip(auto_renew_probability, 0.05, 0.95))

    # Engineered features follow the framework's strict naming rules.
    monthly_active_days_x_product_adoption_count = (
        monthly_active_days * product_adoption_count
    )
    support_tickets_90d_x_support_resolution_days = (
        support_tickets_90d * support_resolution_days
    )
    onboarding_completion_sq = onboarding_completion ** 2

    # The latent log-hazard is constructed from centered / standardized components. This makes the
    # event process directionally coherent while preventing one raw-scale feature from dominating.
    log_relative_hazard = (
        -0.65 * zscore(onboarding_completion)
        -0.35 * zscore(monthly_active_days)
        -0.25 * zscore(product_adoption_count.astype(float))
        +0.20 * zscore(support_tickets_90d.astype(float))
        +0.30 * zscore(support_resolution_days)
        +0.25 * zscore(service_incidents_90d.astype(float))
        -0.15 * zscore(tenure_at_start)
        -0.25 * zscore(auto_renew_flag.astype(float))
        -0.12 * zscore(enterprise_flag.astype(float))
        -0.15 * zscore(monthly_active_days_x_product_adoption_count)
        +0.18 * zscore(support_tickets_90d_x_support_resolution_days)
        -0.10 * zscore(onboarding_completion_sq)
        + rng.normal(0.0, 0.12, n_records)
    )

    # A constant baseline hazard combined with exp(log_relative_hazard) produces data consistent
    # with the proportional-hazards structure used by CoxPH. Censoring is generated independently.
    baseline_monthly_hazard = 0.024
    individual_hazard = baseline_monthly_hazard * np.exp(log_relative_hazard)
    event_time = rng.exponential(scale=1.0 / individual_hazard)
    censor_time = rng.uniform(8.0, max_followup_months, n_records)

    attrition_event = (event_time <= censor_time).astype(int)
    months_observed = np.minimum(event_time, censor_time)
    months_observed = np.clip(months_observed, 0.10, max_followup_months)

    data = pd.DataFrame(
        {
            "customer_id": customer_id,
            "months_observed": months_observed.round(4),
            "attrition_event": attrition_event,
            "onboarding_completion": onboarding_completion.round(6),
            "monthly_active_days": monthly_active_days.round(4),
            "product_adoption_count": product_adoption_count,
            "support_tickets_90d": support_tickets_90d,
            "support_resolution_days": support_resolution_days.round(4),
            "service_incidents_90d": service_incidents_90d,
            "tenure_at_start": tenure_at_start.round(4),
            "auto_renew_flag": auto_renew_flag,
            "enterprise_flag": enterprise_flag,
            # Engineered terms are recalculated below from the stored base-feature values.
            # This avoids false validation differences caused solely by independent rounding.
            "monthly_active_days_x_product_adoption_count": 0.0,
            "support_tickets_90d_x_support_resolution_days": 0.0,
            "onboarding_completion_sq": 0.0,
        }
    )

    # Rebuild engineered terms from the exact persisted base columns. This is the same dependency
    # rule the scenario engine applies after interventions and guarantees a clean baseline audit.
    data["monthly_active_days_x_product_adoption_count"] = (
        data["monthly_active_days"] * data["product_adoption_count"]
    )
    data["support_tickets_90d_x_support_resolution_days"] = (
        data["support_tickets_90d"] * data["support_resolution_days"]
    )
    data["onboarding_completion_sq"] = data["onboarding_completion"] ** 2

    return data


# ==================================================================================================
# CONFIGURATION BUILDERS
# ==================================================================================================

def build_default_retention_config(output_root: Path = Path("outputs")) -> FrameworkConfig:
    return FrameworkConfig(
        project_name="Survival Strategy Framework - Synthetic Retention Demo",
        duration_col="months_observed",
        event_col="attrition_event",
        id_col="customer_id",
        model_features=(
            "onboarding_completion",
            "monthly_active_days",
            "product_adoption_count",
            "support_tickets_90d",
            "support_resolution_days",
            "service_incidents_90d",
            "tenure_at_start",
            "auto_renew_flag",
            "enterprise_flag",
            "monthly_active_days_x_product_adoption_count",
            "support_tickets_90d_x_support_resolution_days",
            "onboarding_completion_sq",
        ),
        segmentation_features=(
            "onboarding_completion",
            "monthly_active_days",
            "product_adoption_count",
            "support_tickets_90d",
            "support_resolution_days",
            "service_incidents_90d",
            "tenure_at_start",
        ),
        profile_features=(
            "onboarding_completion",
            "monthly_active_days",
            "product_adoption_count",
            "support_tickets_90d",
            "support_resolution_days",
            "service_incidents_90d",
            "tenure_at_start",
            "auto_renew_flag",
        ),
        n_clusters=3,
        risk_quantile=0.75,
        penalizer=0.10,
        l1_ratio=0.0,
        random_state=42,
        standardize_model_features=True,
        show_km_confidence_intervals=True,
        run_ph_test=True,
        ph_test_alpha=0.05,
        evaluation_horizon=12.0,
        timeline_points=121,
        output_root=output_root,
        save_input_snapshot=True,
    )


def build_default_retention_scenarios() -> tuple[ScenarioDefinition, ...]:
    """Defines governable, business-interpretable sensitivity scenarios.

    The scenarios alter only the selected high-risk cohort. They are not causal treatment-effect
    estimates. Their purpose is to demonstrate how a governed intervention hypothesis can be
    translated into model-consistent same-cohort sensitivity analysis.
    """
    return (
        ScenarioDefinition(
            name="Onboarding Completion Improvement",
            description=(
                "Increase onboarding completion by 10 percentage points, bounded at 100%, for the "
                "baseline top-quartile risk cohort."
            ),
            changes=(
                FeatureChange(
                    feature="onboarding_completion",
                    operation="add",
                    value=0.10,
                    lower_bound=0.0,
                    upper_bound=1.0,
                    round_digits=6,
                ),
            ),
        ),
        ScenarioDefinition(
            name="Support Resolution Improvement",
            description=(
                "Reduce support resolution time by 25%, with a 0.25-day floor, for the baseline "
                "top-quartile risk cohort."
            ),
            changes=(
                FeatureChange(
                    feature="support_resolution_days",
                    operation="multiply",
                    value=0.75,
                    lower_bound=0.25,
                    upper_bound=20.0,
                    round_digits=4,
                ),
            ),
        ),
        ScenarioDefinition(
            name="Product Adoption Expansion",
            description=(
                "Increase adopted products by one, bounded at eight, for the baseline top-quartile "
                "risk cohort."
            ),
            changes=(
                FeatureChange(
                    feature="product_adoption_count",
                    operation="add",
                    value=1.0,
                    lower_bound=1.0,
                    upper_bound=8.0,
                    round_digits=0,
                ),
            ),
        ),
        ScenarioDefinition(
            name="Combined Retention Strategy",
            description=(
                "Apply the onboarding, support-resolution, and product-adoption assumptions "
                "together to the same baseline top-quartile risk cohort."
            ),
            changes=(
                FeatureChange(
                    feature="onboarding_completion",
                    operation="add",
                    value=0.10,
                    lower_bound=0.0,
                    upper_bound=1.0,
                    round_digits=6,
                ),
                FeatureChange(
                    feature="support_resolution_days",
                    operation="multiply",
                    value=0.75,
                    lower_bound=0.25,
                    upper_bound=20.0,
                    round_digits=4,
                ),
                FeatureChange(
                    feature="product_adoption_count",
                    operation="add",
                    value=1.0,
                    lower_bound=1.0,
                    upper_bound=8.0,
                    round_digits=0,
                ),
            ),
        ),
    )


# ==================================================================================================
# INPUT CONTRACT AND VALIDATION GATEKEEPER
# ==================================================================================================

def engineered_feature_sources(feature: str) -> tuple[str, ...] | None:
    """Returns source columns implied by the strict engineered-feature naming rules."""
    if "_x_" in feature:
        return tuple(part for part in feature.split("_x_") if part)
    if feature.endswith("_sq"):
        return (feature[:-3],)
    return None


def validate_survival_input(
    df: pd.DataFrame,
    config: FrameworkConfig,
    scenarios: Sequence[ScenarioDefinition],
) -> ValidationResult:
    """Validates the full analytical contract before clustering or model fitting begins."""
    result = ValidationResult()

    if not isinstance(df, pd.DataFrame):
        result.add("ERROR", "Dataset", "DATAFRAME", "Input must be a pandas DataFrame.")
        return result
    if df.empty:
        result.add("ERROR", "Dataset", "ROW_COUNT", "Input dataset contains no rows.")
        return result

    required_columns = {
        config.id_col,
        config.duration_col,
        config.event_col,
        *config.model_features,
        *config.segmentation_features,
    }
    missing_columns = sorted(required_columns.difference(df.columns))
    for column in missing_columns:
        result.add("ERROR", "Required Column", column, "Required field is missing from the input.")

    # Configuration checks can run even when input fields are missing.
    if not 0.0 < config.risk_quantile < 1.0:
        result.add("ERROR", "Configuration", "risk_quantile", "Must be between 0 and 1.")
    if config.n_clusters < 2:
        result.add("ERROR", "Configuration", "n_clusters", "Must be at least 2.")
    if config.n_clusters >= len(df):
        result.add("ERROR", "Configuration", "n_clusters", "Must be smaller than row count.")
    if config.penalizer < 0:
        result.add("ERROR", "Configuration", "penalizer", "Must be nonnegative.")
    if not 0.0 <= config.l1_ratio <= 1.0:
        result.add("ERROR", "Configuration", "l1_ratio", "Must be between 0 and 1.")
    if config.evaluation_horizon <= 0:
        result.add("ERROR", "Configuration", "evaluation_horizon", "Must be positive.")
    if config.timeline_points < 20:
        result.add("ERROR", "Configuration", "timeline_points", "Must be at least 20.")

    if missing_columns:
        return result

    # Governed ID validation.
    if df[config.id_col].isna().any():
        result.add("ERROR", "Identifier", config.id_col, "ID contains missing values.")
    if df[config.id_col].duplicated().any():
        duplicate_count = int(df[config.id_col].duplicated(keep=False).sum())
        result.add(
            "ERROR",
            "Identifier",
            config.id_col,
            f"ID is not unique; {duplicate_count:,} rows participate in duplicates.",
        )

    # Duration validation.
    duration = pd.to_numeric(df[config.duration_col], errors="coerce")
    if duration.isna().any():
        result.add("ERROR", "Time-to-Event", config.duration_col, "Duration must be numeric and non-missing.")
    else:
        if not np.isfinite(duration.to_numpy()).all():
            result.add("ERROR", "Time-to-Event", config.duration_col, "Duration contains infinite values.")
        if (duration <= 0).any():
            result.add("ERROR", "Time-to-Event", config.duration_col, "Duration must be strictly positive.")
        if config.evaluation_horizon > float(duration.max()):
            result.add(
                "WARNING",
                "Time-to-Event",
                "evaluation_horizon",
                "Evaluation horizon exceeds the maximum observed duration; predicted survival is an extrapolative model output.",
            )

    # Event validation.
    event = pd.to_numeric(df[config.event_col], errors="coerce")
    if event.isna().any():
        result.add("ERROR", "Event", config.event_col, "Event field must be binary and non-missing.")
    else:
        invalid_events = sorted(set(event.unique()).difference({0, 1}))
        if invalid_events:
            result.add(
                "ERROR",
                "Event",
                config.event_col,
                f"Event field contains values outside {{0, 1}}: {invalid_events[:10]}",
            )
        else:
            event_rate = float(event.mean())
            result.add(
                "INFO",
                "Event",
                config.event_col,
                f"Observed event rate is {event_rate:.2%}; right-censoring rate is {1.0 - event_rate:.2%}.",
            )
            if event_rate < config.low_event_rate_warning:
                result.add("WARNING", "Event", config.event_col, "Event rate is very low; model stability requires review.")
            if event_rate > config.high_event_rate_warning:
                result.add("WARNING", "Event", config.event_col, "Event rate is very high; censoring design requires review.")

    # Feature validation. Version 1 of the hardened framework requires numeric / Boolean inputs.
    all_features = list(dict.fromkeys((*config.model_features, *config.segmentation_features)))
    for feature in all_features:
        series = df[feature]
        if not (pd.api.types.is_numeric_dtype(series) or pd.api.types.is_bool_dtype(series)):
            result.add(
                "ERROR",
                "Feature Contract",
                feature,
                "Feature must be numeric or Boolean; automatic categorical encoding is not claimed in this release.",
            )
            continue

        numeric = pd.to_numeric(series, errors="coerce")
        missing_ratio = float(numeric.isna().mean())
        if missing_ratio > 0:
            result.add(
                "ERROR",
                "Feature Contract",
                feature,
                f"Feature contains {missing_ratio:.2%} missing/non-numeric values; no imputation rule is configured.",
            )
        if missing_ratio >= config.high_missingness_warning:
            result.add("WARNING", "Feature Contract", feature, "Missingness exceeds the configured warning threshold.")
        finite = numeric.dropna().to_numpy()
        if finite.size and not np.isfinite(finite).all():
            result.add("ERROR", "Feature Contract", feature, "Feature contains infinite values.")
        if numeric.nunique(dropna=True) <= 1:
            severity = "ERROR" if feature in config.model_features else "WARNING"
            result.add(severity, "Feature Contract", feature, "Feature has zero variance.")

    # Engineered-feature validation confirms that stored interactions / squares agree with base data.
    for feature in config.model_features:
        sources = engineered_feature_sources(feature)
        if not sources:
            continue
        missing_sources = [source for source in sources if source not in df.columns]
        if missing_sources:
            result.add(
                "ERROR",
                "Engineered Feature",
                feature,
                f"Source field(s) missing: {missing_sources}",
            )
            continue

        if "_x_" in feature:
            expected = df.loc[:, list(sources)].prod(axis=1)
        else:
            expected = pd.to_numeric(df[sources[0]], errors="coerce") ** 2

        actual = pd.to_numeric(df[feature], errors="coerce")
        max_difference = float(np.nanmax(np.abs(actual.to_numpy() - expected.to_numpy())))
        if not np.isfinite(max_difference) or max_difference > 1e-6:
            result.add(
                "ERROR",
                "Engineered Feature",
                feature,
                f"Stored term is out of synchronization with its source field(s); max absolute difference={max_difference:.6g}.",
            )

    # Scenario controls.
    scenario_names = [scenario.name for scenario in scenarios]
    duplicate_names = sorted({name for name in scenario_names if scenario_names.count(name) > 1})
    for name in duplicate_names:
        result.add("ERROR", "Scenario", name, "Scenario names must be unique.")

    protected_fields = {config.id_col, config.duration_col, config.event_col}
    engineered_fields = {
        feature for feature in config.model_features if engineered_feature_sources(feature)
    }
    for scenario in scenarios:
        if not scenario.changes:
            result.add("WARNING", "Scenario", scenario.name, "Scenario contains no feature changes.")
        for change in scenario.changes:
            if change.feature not in df.columns:
                result.add("ERROR", "Scenario", scenario.name, f"Scenario feature '{change.feature}' does not exist.")
            if change.feature in protected_fields:
                result.add("ERROR", "Scenario", scenario.name, f"Protected field '{change.feature}' may not be changed.")
            if change.feature in engineered_fields:
                result.add(
                    "ERROR",
                    "Scenario",
                    scenario.name,
                    f"Engineered field '{change.feature}' may not be changed directly; change its base feature(s).",
                )
            if change.operation not in {"multiply", "add", "replace"}:
                result.add("ERROR", "Scenario", scenario.name, f"Unsupported operation '{change.operation}'.")
            if (
                change.lower_bound is not None
                and change.upper_bound is not None
                and change.lower_bound > change.upper_bound
            ):
                result.add("ERROR", "Scenario", scenario.name, "Lower bound exceeds upper bound.")

    return result


# ==================================================================================================
# PERSONA BRANCH — UNSUPERVISED SEGMENTATION
# ==================================================================================================

def fit_personas(df: pd.DataFrame, config: FrameworkConfig) -> PersonaArtifacts:
    """Fits reproducible K-Means personas independently from the CoxPH model."""
    features = numeric_frame(df, config.segmentation_features)
    scaler = StandardScaler()
    scaled = scaler.fit_transform(features)

    kmeans = KMeans(
        n_clusters=config.n_clusters,
        random_state=config.random_state,
        n_init=20,
    )
    raw_labels = kmeans.fit_predict(scaled)

    scored = df.copy()
    scored["Segment_ID"] = raw_labels.astype(int)
    scored["Persona"] = scored["Segment_ID"].map(
        lambda segment: f"Persona Group {int(segment) + 1}"
    )

    grouped = scored.groupby("Persona", sort=True)
    profiles = grouped.agg(
        Record_Count=(config.id_col, "count"),
        Event_Rate=(config.event_col, "mean"),
        Median_Duration=(config.duration_col, "median"),
    ).reset_index()

    profile_features = config.profile_features or config.segmentation_features
    feature_means = grouped[list(profile_features)].mean().reset_index()
    profiles = profiles.merge(feature_means, on="Persona", how="left")

    return PersonaArtifacts(
        scored_data=scored,
        scaler=scaler,
        kmeans=kmeans,
        profiles=profiles,
    )


# ==================================================================================================
# COXPH BRANCH — MODEL FITTING, HAZARD RATIOS, AND ASSUMPTION TESTING
# ==================================================================================================

def transform_model_features(
    df: pd.DataFrame,
    config: FrameworkConfig,
    scaler: StandardScaler | None,
) -> pd.DataFrame:
    raw = numeric_frame(df, config.model_features)
    if scaler is None:
        return raw
    transformed = scaler.transform(raw)
    return pd.DataFrame(transformed, columns=list(config.model_features), index=df.index)


def _model_summary_frame(cph: CoxPHFitter, standardized: bool) -> pd.DataFrame:
    summary = cph.summary.reset_index()
    first_column = summary.columns[0]
    summary = summary.rename(columns={first_column: "Feature"})

    rename_map = {
        "coef": "Coefficient",
        "exp(coef)": "Hazard_Ratio",
        "se(coef)": "Standard_Error",
        "coef lower 95%": "Coefficient_CI_Lower",
        "coef upper 95%": "Coefficient_CI_Upper",
        "exp(coef) lower 95%": "Hazard_Ratio_CI_Lower",
        "exp(coef) upper 95%": "Hazard_Ratio_CI_Upper",
        "z": "Z_Statistic",
        "p": "P_Value",
        "-log2(p)": "Negative_Log2_P",
    }
    summary = summary.rename(columns=rename_map)

    desired = [
        "Feature",
        "Coefficient",
        "Hazard_Ratio",
        "Standard_Error",
        "Coefficient_CI_Lower",
        "Coefficient_CI_Upper",
        "Hazard_Ratio_CI_Lower",
        "Hazard_Ratio_CI_Upper",
        "Z_Statistic",
        "P_Value",
        "Negative_Log2_P",
    ]
    available = [column for column in desired if column in summary.columns]
    summary = summary.loc[:, available].copy()
    summary["Interpretation_Unit"] = "Per 1 SD" if standardized else "Per raw feature unit"
    return summary


def fit_cox_model(df: pd.DataFrame, config: FrameworkConfig) -> CoxModelArtifacts:
    """Fits the regularized CoxPH branch and captures technical evidence.

    The persona fields are never included here. The model matrix is built only from the governed
    model feature list. If model standardization is enabled, one scaler is fitted on baseline data
    and reused unchanged for all scenario predictions.
    """
    raw_features = numeric_frame(df, config.model_features)
    model_scaler: StandardScaler | None = None

    if config.standardize_model_features:
        model_scaler = StandardScaler()
        transformed_array = model_scaler.fit_transform(raw_features)
        transformed_features = pd.DataFrame(
            transformed_array,
            columns=list(config.model_features),
            index=df.index,
        )
    else:
        transformed_features = raw_features

    training_frame = transformed_features.copy()
    training_frame[config.duration_col] = pd.to_numeric(df[config.duration_col], errors="raise")
    training_frame[config.event_col] = pd.to_numeric(df[config.event_col], errors="raise").astype(int)

    cph = CoxPHFitter(penalizer=config.penalizer, l1_ratio=config.l1_ratio)
    captured_warnings: list[str] = []

    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        cph.fit(
            training_frame,
            duration_col=config.duration_col,
            event_col=config.event_col,
            show_progress=False,
        )
        captured_warnings = [str(item.message) for item in caught]

    model_summary = _model_summary_frame(cph, config.standardize_model_features)

    fit_metadata = pd.DataFrame(
        [
            {
                "Engine_Release": ENGINE_RELEASE,
                "Rows": int(len(training_frame)),
                "Events": int(training_frame[config.event_col].sum()),
                "Censored": int(len(training_frame) - training_frame[config.event_col].sum()),
                "Event_Rate": float(training_frame[config.event_col].mean()),
                "Feature_Count": int(len(config.model_features)),
                "Penalizer": float(config.penalizer),
                "L1_Ratio": float(config.l1_ratio),
                "Standardized_Features": bool(config.standardize_model_features),
                "Concordance_Index": float(cph.concordance_index_),
                "Partial_AIC": float(cph.AIC_partial_),
                "Log_Likelihood": float(cph.log_likelihood_),
                "Fit_Status": "SUCCESS",
            }
        ]
    )

    if config.run_ph_test:
        try:
            ph_result = proportional_hazard_test(
                cph,
                training_frame,
                time_transform="rank",
            )
            ph_test_results = ph_result.summary.reset_index()
            first_column = ph_test_results.columns[0]
            ph_test_results = ph_test_results.rename(
                columns={
                    first_column: "Feature",
                    "test_statistic": "Test_Statistic",
                    "p": "P_Value",
                }
            )
            ph_test_results["PH_Assumption_Status"] = np.where(
                ph_test_results["P_Value"] < config.ph_test_alpha,
                "REVIEW",
                "PASS",
            )
            ph_test_results["Alpha"] = config.ph_test_alpha
        except Exception as exc:  # PH testing is evidence; failure is recorded rather than hidden.
            ph_test_results = pd.DataFrame(
                [
                    {
                        "Feature": "ALL",
                        "Test_Statistic": np.nan,
                        "P_Value": np.nan,
                        "PH_Assumption_Status": "TEST_ERROR",
                        "Alpha": config.ph_test_alpha,
                        "Error": str(exc),
                    }
                ]
            )
    else:
        ph_test_results = pd.DataFrame(
            [
                {
                    "Feature": "ALL",
                    "Test_Statistic": np.nan,
                    "P_Value": np.nan,
                    "PH_Assumption_Status": "NOT_RUN",
                    "Alpha": config.ph_test_alpha,
                }
            ]
        )

    fit_warnings = pd.DataFrame(
        [{"Warning": message} for message in captured_warnings]
        or [{"Warning": "No fit warnings were captured."}]
    )

    return CoxModelArtifacts(
        model=cph,
        model_scaler=model_scaler,
        training_frame=training_frame,
        model_summary=model_summary,
        fit_metadata=fit_metadata,
        ph_test_results=ph_test_results,
        fit_warnings=fit_warnings,
    )


# ==================================================================================================
# POPULATION RISK SCORING AND GOVERNED TARGET-COHORT SELECTION
# ==================================================================================================

def score_and_stratify_population(
    df: pd.DataFrame,
    personas: PersonaArtifacts,
    model_artifacts: CoxModelArtifacts,
    config: FrameworkConfig,
) -> RiskArtifacts:
    model_matrix = transform_model_features(
        df,
        config,
        model_artifacts.model_scaler,
    )
    partial_hazard = model_artifacts.model.predict_partial_hazard(model_matrix)
    hazard_series = pd.Series(
        np.asarray(partial_hazard).reshape(-1),
        index=df.index,
        name="Partial_Hazard",
    )

    scored = df.copy()
    scored["Persona"] = personas.scored_data["Persona"]
    scored["Segment_ID"] = personas.scored_data["Segment_ID"]
    scored["Partial_Hazard"] = hazard_series
    scored["Risk_Percentile"] = scored["Partial_Hazard"].rank(
        method="average",
        pct=True,
    )

    q25 = float(scored["Partial_Hazard"].quantile(0.25))
    target_threshold = float(scored["Partial_Hazard"].quantile(config.risk_quantile))
    target_label = f"High Risk (Top {1.0 - config.risk_quantile:.0%})"

    scored["Risk_Tier"] = "Mid Risk"
    scored.loc[scored["Partial_Hazard"] <= q25, "Risk_Tier"] = "Benchmark (Bottom 25%)"
    scored.loc[
        scored["Partial_Hazard"] >= target_threshold,
        "Risk_Tier",
    ] = target_label
    scored["Target_Cohort_Flag"] = (
        scored["Partial_Hazard"] >= target_threshold
    ).astype(int)

    profile_features = config.profile_features or config.segmentation_features
    grouped = scored.groupby("Risk_Tier", sort=False)
    risk_profiles = grouped.agg(
        Record_Count=(config.id_col, "count"),
        Event_Rate=(config.event_col, "mean"),
        Median_Duration=(config.duration_col, "median"),
        Mean_Partial_Hazard=("Partial_Hazard", "mean"),
    ).reset_index()
    feature_means = grouped[list(profile_features)].mean().reset_index()
    risk_profiles = risk_profiles.merge(feature_means, on="Risk_Tier", how="left")

    target_index = scored.index[scored["Target_Cohort_Flag"] == 1]
    target_columns = [
        config.id_col,
        "Persona",
        "Risk_Tier",
        "Partial_Hazard",
        "Risk_Percentile",
        *[feature for feature in profile_features if feature in scored.columns],
    ]
    target_cohort = scored.loc[target_index, list(dict.fromkeys(target_columns))].copy()

    return RiskArtifacts(
        scored_data=scored,
        risk_tier_profiles=risk_profiles,
        target_cohort=target_cohort,
        target_index=target_index,
        target_threshold=target_threshold,
        target_label=target_label,
    )


# ==================================================================================================
# ENGINEERED-FEATURE DEPENDENCY SYNCHRONIZATION
# ==================================================================================================

def sync_engineered_dependencies(
    df: pd.DataFrame,
    model_features: Sequence[str],
    scenario_name: str,
) -> tuple[pd.DataFrame, pd.DataFrame]:
    """Rebuilds interaction and squared terms after a scenario changes base features.

    The audit records both the pre-rebuild discrepancy and the post-rebuild verification. This
    turns dependency synchronization from an undocumented convenience into explicit control
    evidence.
    """
    synced = df.copy()
    audit_rows: list[dict[str, Any]] = []

    for feature in model_features:
        sources = engineered_feature_sources(feature)
        if not sources:
            continue

        if any(source not in synced.columns for source in sources):
            audit_rows.append(
                {
                    "Scenario_Name": scenario_name,
                    "Engineered_Feature": feature,
                    "Dependency_Type": "INTERACTION" if "_x_" in feature else "SQUARE",
                    "Source_Features": " | ".join(sources),
                    "Pre_Rebuild_Max_Abs_Diff": np.nan,
                    "Post_Rebuild_Max_Abs_Diff": np.nan,
                    "Recalculated": False,
                    "Validation_Status": "FAIL_MISSING_SOURCE",
                }
            )
            continue

        previous = pd.to_numeric(synced[feature], errors="coerce")
        if "_x_" in feature:
            expected = synced.loc[:, list(sources)].prod(axis=1)
            dependency_type = "INTERACTION"
        else:
            expected = pd.to_numeric(synced[sources[0]], errors="coerce") ** 2
            dependency_type = "SQUARE"

        pre_difference = float(np.nanmax(np.abs(previous.to_numpy() - expected.to_numpy())))
        synced[feature] = expected
        post_difference = float(
            np.nanmax(
                np.abs(
                    pd.to_numeric(synced[feature], errors="coerce").to_numpy()
                    - expected.to_numpy()
                )
            )
        )

        audit_rows.append(
            {
                "Scenario_Name": scenario_name,
                "Engineered_Feature": feature,
                "Dependency_Type": dependency_type,
                "Source_Features": " | ".join(sources),
                "Pre_Rebuild_Max_Abs_Diff": pre_difference,
                "Post_Rebuild_Max_Abs_Diff": post_difference,
                "Recalculated": True,
                "Validation_Status": "PASS" if post_difference <= 1e-10 else "FAIL",
            }
        )

    audit = pd.DataFrame(audit_rows)
    if audit.empty:
        audit = pd.DataFrame(
            [
                {
                    "Scenario_Name": scenario_name,
                    "Engineered_Feature": "NONE",
                    "Dependency_Type": "NONE",
                    "Source_Features": "",
                    "Pre_Rebuild_Max_Abs_Diff": 0.0,
                    "Post_Rebuild_Max_Abs_Diff": 0.0,
                    "Recalculated": False,
                    "Validation_Status": "NOT_APPLICABLE",
                }
            ]
        )
    return synced, audit


# ==================================================================================================
# GENERIC SAME-COHORT SCENARIO SIMULATION
# ==================================================================================================

def apply_feature_change(
    df: pd.DataFrame,
    row_index: pd.Index,
    change: FeatureChange,
) -> pd.DataFrame:
    """Applies one bounded change to the governed target cohort only."""
    result = df.copy()
    values = pd.to_numeric(result.loc[row_index, change.feature], errors="raise").astype(float)

    if change.operation == "multiply":
        changed = values * change.value
    elif change.operation == "add":
        changed = values + change.value
    elif change.operation == "replace":
        changed = pd.Series(change.value, index=values.index, dtype=float)
    else:  # Protected by validation, retained here as a fail-fast defensive check.
        raise ValueError(f"Unsupported scenario operation: {change.operation}")

    if change.lower_bound is not None:
        changed = changed.clip(lower=change.lower_bound)
    if change.upper_bound is not None:
        changed = changed.clip(upper=change.upper_bound)
    if change.round_digits is not None:
        changed = changed.round(change.round_digits)

    result.loc[row_index, change.feature] = changed
    return result


def mean_predicted_survival(
    model: CoxPHFitter,
    model_matrix: pd.DataFrame,
    times: np.ndarray,
) -> pd.Series:
    """Returns the cohort-average CoxPH-predicted survival function."""
    survival = model.predict_survival_function(model_matrix, times=times)
    return survival.mean(axis=1)


def simulate_scenarios(
    source_df: pd.DataFrame,
    config: FrameworkConfig,
    scenarios: Sequence[ScenarioDefinition],
    model_artifacts: CoxModelArtifacts,
    risk_artifacts: RiskArtifacts,
) -> ScenarioArtifacts:
    """Applies controlled scenario changes and re-scores the identical target cohort."""
    target_index = risk_artifacts.target_index
    if len(target_index) == 0:
        raise ValueError("Target cohort is empty; revise risk_quantile or input data.")

    baseline_matrix_all = transform_model_features(
        source_df,
        config,
        model_artifacts.model_scaler,
    )
    baseline_matrix = baseline_matrix_all.loc[target_index]

    baseline_hazard = pd.Series(
        np.asarray(model_artifacts.model.predict_partial_hazard(baseline_matrix)).reshape(-1),
        index=target_index,
        name="Baseline_Partial_Hazard",
    )
    baseline_mean_hazard = float(baseline_hazard.mean())

    max_time = max(float(source_df[config.duration_col].max()), config.evaluation_horizon)
    timeline = np.linspace(0.0, max_time, config.timeline_points)
    baseline_curve = mean_predicted_survival(
        model_artifacts.model,
        baseline_matrix,
        timeline,
    )
    baseline_horizon = float(
        mean_predicted_survival(
            model_artifacts.model,
            baseline_matrix,
            np.array([config.evaluation_horizon]),
        ).iloc[0]
    )

    summary_rows: list[dict[str, Any]] = []
    score_frames: list[pd.DataFrame] = []
    dependency_frames: list[pd.DataFrame] = []
    curve_frame = pd.DataFrame(
        {
            "Time": timeline,
            "Baseline Target Cohort": baseline_curve.to_numpy(),
        }
    )

    for scenario in scenarios:
        scenario_df = source_df.copy()
        for change in scenario.changes:
            scenario_df = apply_feature_change(
                scenario_df,
                target_index,
                change,
            )

        scenario_df, dependency_audit = sync_engineered_dependencies(
            scenario_df,
            config.model_features,
            scenario.name,
        )
        dependency_frames.append(dependency_audit)

        # Re-run the strict contract after transformation. Scenarios may not create missing,
        # infinite, out-of-range, or stale engineered values.
        scenario_validation = validate_survival_input(scenario_df, config, scenarios=())
        scenario_validation.raise_for_errors()

        scenario_matrix_all = transform_model_features(
            scenario_df,
            config,
            model_artifacts.model_scaler,
        )
        scenario_matrix = scenario_matrix_all.loc[target_index]

        scenario_hazard = pd.Series(
            np.asarray(model_artifacts.model.predict_partial_hazard(scenario_matrix)).reshape(-1),
            index=target_index,
            name="Scenario_Partial_Hazard",
        )
        scenario_mean_hazard = float(scenario_hazard.mean())
        relative_hazard_reduction = (
            (baseline_mean_hazard - scenario_mean_hazard) / baseline_mean_hazard
            if baseline_mean_hazard != 0
            else np.nan
        )

        scenario_curve = mean_predicted_survival(
            model_artifacts.model,
            scenario_matrix,
            timeline,
        )
        curve_frame[scenario.name] = scenario_curve.to_numpy()

        scenario_horizon = float(
            mean_predicted_survival(
                model_artifacts.model,
                scenario_matrix,
                np.array([config.evaluation_horizon]),
            ).iloc[0]
        )
        survival_uplift = scenario_horizon - baseline_horizon

        summary_rows.append(
            {
                "Scenario_Name": scenario.name,
                "Description": scenario.description,
                "Target_Cohort_Count": int(len(target_index)),
                "Baseline_Mean_Partial_Hazard": baseline_mean_hazard,
                "Scenario_Mean_Partial_Hazard": scenario_mean_hazard,
                "Relative_Hazard_Reduction": relative_hazard_reduction,
                "Evaluation_Horizon": float(config.evaluation_horizon),
                "Baseline_Survival_At_Horizon": baseline_horizon,
                "Scenario_Survival_At_Horizon": scenario_horizon,
                "Survival_Probability_Uplift": survival_uplift,
                "Modeled_Direction": (
                    "IMPROVED"
                    if relative_hazard_reduction > 0 and survival_uplift > 0
                    else "ADVERSE_OR_MIXED"
                ),
            }
        )

        score_frame = pd.DataFrame(
            {
                config.id_col: source_df.loc[target_index, config.id_col].to_numpy(),
                "Scenario_Name": scenario.name,
                "Baseline_Partial_Hazard": baseline_hazard.to_numpy(),
                "Scenario_Partial_Hazard": scenario_hazard.to_numpy(),
                "Partial_Hazard_Delta": (
                    scenario_hazard.to_numpy() - baseline_hazard.to_numpy()
                ),
            },
            index=target_index,
        )
        score_frame["Relative_Partial_Hazard_Change"] = np.where(
            score_frame["Baseline_Partial_Hazard"] != 0,
            score_frame["Partial_Hazard_Delta"]
            / score_frame["Baseline_Partial_Hazard"],
            np.nan,
        )
        score_frames.append(score_frame.reset_index(drop=True))

    summary = pd.DataFrame(summary_rows)
    if not summary.empty:
        summary = summary.sort_values(
            ["Relative_Hazard_Reduction", "Survival_Probability_Uplift"],
            ascending=[False, False],
        ).reset_index(drop=True)
        summary["Scenario_Rank"] = np.arange(1, len(summary) + 1)

    target_scores = (
        pd.concat(score_frames, ignore_index=True)
        if score_frames
        else pd.DataFrame()
    )
    dependency_audit = (
        pd.concat(dependency_frames, ignore_index=True)
        if dependency_frames
        else pd.DataFrame()
    )

    return ScenarioArtifacts(
        summary=summary,
        target_scores=target_scores,
        dependency_audit=dependency_audit,
        predicted_survival_curves=curve_frame,
    )


# ==================================================================================================
# EVIDENCE VISUALS
# ==================================================================================================

def _prepare_dark_axis(ax: plt.Axes) -> None:
    ax.set_facecolor(COLOR_BACKGROUND)
    ax.tick_params(colors=COLOR_LIGHT_GRAY)
    for spine in ax.spines.values():
        spine.set_color(COLOR_MID_GRAY)
    ax.xaxis.label.set_color(COLOR_LIGHT_GRAY)
    ax.yaxis.label.set_color(COLOR_LIGHT_GRAY)
    ax.title.set_color(COLOR_WHITE)
    ax.grid(alpha=0.15, color=COLOR_MID_GRAY)


def plot_persona_kaplan_meier(
    data: pd.DataFrame,
    config: FrameworkConfig,
    output_path: Path,
) -> None:
    figure, axis = plt.subplots(figsize=(10, 6), facecolor=COLOR_BACKGROUND)
    _prepare_dark_axis(axis)
    kmf = KaplanMeierFitter()

    for persona in sorted(data["Persona"].unique()):
        subset = data[data["Persona"] == persona]
        kmf.fit(
            subset[config.duration_col],
            event_observed=subset[config.event_col],
            label=persona,
        )
        kmf.plot_survival_function(
            ax=axis,
            ci_show=config.show_km_confidence_intervals,
        )

    axis.set_title("Persona Kaplan-Meier Survival Evidence")
    axis.set_xlabel(f"Time ({config.duration_col})")
    axis.set_ylabel("Observed Probability of Remaining Event-Free")
    legend = axis.legend(facecolor=COLOR_PANEL, edgecolor=COLOR_MID_GRAY)
    for text in legend.get_texts():
        text.set_color(COLOR_WHITE)
    figure.tight_layout()
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def plot_risk_tier_kaplan_meier(
    data: pd.DataFrame,
    config: FrameworkConfig,
    output_path: Path,
) -> None:
    figure, axis = plt.subplots(figsize=(10, 6), facecolor=COLOR_BACKGROUND)
    _prepare_dark_axis(axis)
    kmf = KaplanMeierFitter()

    tier_order = [
        f"High Risk (Top {1.0 - config.risk_quantile:.0%})",
        "Mid Risk",
        "Benchmark (Bottom 25%)",
    ]
    for tier in tier_order:
        subset = data[data["Risk_Tier"] == tier]
        if subset.empty:
            continue
        kmf.fit(
            subset[config.duration_col],
            event_observed=subset[config.event_col],
            label=tier,
        )
        kmf.plot_survival_function(
            ax=axis,
            ci_show=config.show_km_confidence_intervals,
        )

    axis.set_title("Observed Kaplan-Meier Survival by Model-Derived Risk Tier")
    axis.set_xlabel(f"Time ({config.duration_col})")
    axis.set_ylabel("Observed Probability of Remaining Event-Free")
    legend = axis.legend(facecolor=COLOR_PANEL, edgecolor=COLOR_MID_GRAY)
    for text in legend.get_texts():
        text.set_color(COLOR_WHITE)
    figure.tight_layout()
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def plot_predicted_survival_comparison(
    curves: pd.DataFrame,
    scenario_summary: pd.DataFrame,
    output_path: Path,
) -> None:
    figure, axis = plt.subplots(figsize=(10, 6), facecolor=COLOR_BACKGROUND)
    _prepare_dark_axis(axis)

    axis.plot(
        curves["Time"],
        curves["Baseline Target Cohort"],
        label="Baseline Target Cohort",
        linewidth=3.0,
        color=COLOR_LIGHT_GRAY,
    )

    palette = [COLOR_CYAN, COLOR_GOLD, COLOR_TEAL, COLOR_GREEN, COLOR_RED]
    ordered_scenarios = (
        scenario_summary.sort_values("Scenario_Rank")["Scenario_Name"].tolist()
        if not scenario_summary.empty
        else []
    )
    for index, scenario_name in enumerate(ordered_scenarios):
        if scenario_name not in curves.columns:
            continue
        axis.plot(
            curves["Time"],
            curves[scenario_name],
            label=scenario_name,
            linewidth=2.2,
            color=palette[index % len(palette)],
        )

    axis.set_title("CoxPH-Predicted Survival: Same Target Cohort, Controlled Scenarios")
    axis.set_xlabel("Time")
    axis.set_ylabel("Mean Predicted Survival Probability")
    axis.set_ylim(0.0, 1.02)
    legend = axis.legend(facecolor=COLOR_PANEL, edgecolor=COLOR_MID_GRAY, fontsize=8)
    for text in legend.get_texts():
        text.set_color(COLOR_WHITE)
    figure.tight_layout()
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def plot_scenario_hazard_reduction(
    summary: pd.DataFrame,
    output_path: Path,
) -> None:
    figure, axis = plt.subplots(figsize=(10, 5.5), facecolor=COLOR_BACKGROUND)
    _prepare_dark_axis(axis)

    if summary.empty:
        axis.text(0.5, 0.5, "No configured scenarios", ha="center", va="center", color=COLOR_WHITE)
    else:
        ordered = summary.sort_values("Relative_Hazard_Reduction", ascending=True)
        values = ordered["Relative_Hazard_Reduction"] * 100.0
        axis.barh(ordered["Scenario_Name"], values, color=COLOR_CYAN, alpha=0.85)
        axis.axvline(0.0, color=COLOR_LIGHT_GRAY, linewidth=1)
        for position, value in enumerate(values):
            axis.text(
                value + (0.25 if value >= 0 else -0.25),
                position,
                f"{value:.1f}%",
                va="center",
                ha="left" if value >= 0 else "right",
                color=COLOR_WHITE,
                fontsize=9,
            )

    axis.set_title("Modeled Relative Hazard Reduction by Scenario")
    axis.set_xlabel("Relative Change in Mean Partial Hazard (%)")
    figure.tight_layout()
    figure.savefig(output_path, dpi=180, bbox_inches="tight", facecolor=figure.get_facecolor())
    plt.close(figure)


def generate_evidence_plots(
    run_directory: Path,
    config: FrameworkConfig,
    personas: PersonaArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
) -> PlotArtifacts:
    chart_directory = run_directory / "charts"
    persona_path = chart_directory / "persona_kaplan_meier.png"
    risk_path = chart_directory / "risk_tier_kaplan_meier.png"
    scenario_path = chart_directory / "baseline_vs_scenario_survival.png"
    comparison_path = chart_directory / "scenario_hazard_reduction.png"

    plot_persona_kaplan_meier(personas.scored_data, config, persona_path)
    plot_risk_tier_kaplan_meier(risk.scored_data, config, risk_path)
    plot_predicted_survival_comparison(
        scenarios.predicted_survival_curves,
        scenarios.summary,
        scenario_path,
    )
    plot_scenario_hazard_reduction(scenarios.summary, comparison_path)

    return PlotArtifacts(
        persona_km=persona_path,
        risk_tier_km=risk_path,
        baseline_vs_scenario=scenario_path,
        scenario_comparison=comparison_path,
    )


# ==================================================================================================
# EXECUTIVE INTERPRETATION
# ==================================================================================================

def build_executive_narrative(
    config: FrameworkConfig,
    model: CoxModelArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
) -> str:
    event_rate = float(model.fit_metadata.iloc[0]["Event_Rate"])
    concordance = float(model.fit_metadata.iloc[0]["Concordance_Index"])
    target_count = int(len(risk.target_cohort))

    if scenarios.summary.empty:
        scenario_text = "No intervention scenarios were configured for this run."
    else:
        best = scenarios.summary.sort_values("Scenario_Rank").iloc[0]
        scenario_text = (
            f"Among the configured illustrative scenarios, '{best['Scenario_Name']}' produced "
            f"the largest modeled improvement for the same {target_count:,}-record target cohort: "
            f"mean partial hazard declined by {best['Relative_Hazard_Reduction']:.2%}, and modeled "
            f"survival probability at {config.evaluation_horizon:g} time units increased by "
            f"{best['Survival_Probability_Uplift']:.2%}."
        )

    return (
        f"The {config.project_name} run analyzed {len(risk.scored_data):,} records with an observed "
        f"event rate of {event_rate:.2%}. The regularized CoxPH model produced a concordance index "
        f"of {concordance:.3f}. The baseline {risk.target_label.lower()} rule selected {target_count:,} records "
        f"for same-cohort sensitivity testing. {scenario_text} These outputs are model-based "
        f"sensitivity evidence from demonstration data; they are not causal treatment effects, "
        f"production retention forecasts, or realized financial outcomes."
    )


# ==================================================================================================
# EXECUTIVE POWERPOINT
# ==================================================================================================

def _ppt_set_background(slide: Any, rgb: RGBColor = RGBColor(7, 19, 26)) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _ppt_add_title(slide: Any, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.size = Pt(24)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(244, 245, 246)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.85), Inches(12.0), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].font.name = "Arial"
        subtitle_frame.paragraphs[0].font.size = Pt(12)
        subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(32, 215, 229)


def _ppt_add_footer(slide: Any) -> None:
    footer_box = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.25))
    frame = footer_box.text_frame
    frame.text = "Public demonstration data | Illustrative scenarios | Not a production retention forecast"
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.name = "Arial"
    paragraph.font.size = Pt(8)
    paragraph.font.color.rgb = RGBColor(113, 128, 138)


def _ppt_add_picture_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    subtitle: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_background(slide)
    _ppt_add_title(slide, title, subtitle)
    slide.shapes.add_picture(str(image_path), Inches(1.15), Inches(1.30), width=Inches(11.0))
    _ppt_add_footer(slide)


def _ppt_add_scenario_table(slide: Any, summary: pd.DataFrame) -> None:
    display_columns = [
        "Scenario_Name",
        "Relative_Hazard_Reduction",
        "Survival_Probability_Uplift",
        "Scenario_Rank",
    ]
    table_df = summary.loc[:, display_columns].copy()
    table_df["Relative_Hazard_Reduction"] = table_df["Relative_Hazard_Reduction"].map(
        lambda value: f"{value:.1%}"
    )
    table_df["Survival_Probability_Uplift"] = table_df["Survival_Probability_Uplift"].map(
        lambda value: f"{value:.1%}"
    )
    table_df = table_df.rename(
        columns={
            "Scenario_Name": "Scenario",
            "Relative_Hazard_Reduction": "Hazard Reduction",
            "Survival_Probability_Uplift": "Survival Uplift",
            "Scenario_Rank": "Rank",
        }
    )

    rows, columns = len(table_df) + 1, len(table_df.columns)
    shape = slide.shapes.add_table(
        rows,
        columns,
        Inches(0.70),
        Inches(1.40),
        Inches(12.0),
        Inches(2.7),
    )
    table = shape.table
    for col_idx, column in enumerate(table_df.columns):
        table.cell(0, col_idx).text = column
    for row_idx, (_, row) in enumerate(table_df.iterrows(), start=1):
        for col_idx, value in enumerate(row.tolist()):
            table.cell(row_idx, col_idx).text = str(value)

    for row_idx in range(rows):
        for col_idx in range(columns):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(12, 32, 42) if row_idx else RGBColor(20, 73, 88)
            )
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Arial"
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = RGBColor(244, 245, 246)
                if row_idx == 0:
                    paragraph.font.bold = True


def export_executive_powerpoint(
    output_path: Path,
    config: FrameworkConfig,
    model: CoxModelArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
    plots: PlotArtifacts,
    narrative: str,
) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — Purpose and boundary.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_background(slide)
    _ppt_add_title(
        slide,
        "Survival Strategy Framework",
        "From time-to-event data to governed intervention evidence",
    )
    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.6), Inches(3.8))
    text_frame = body.text_frame
    text_frame.word_wrap = True
    bullets = [
        "Separate persona discovery from CoxPH risk modeling.",
        "Select a model-derived top-quartile target cohort.",
        "Apply controlled scenarios to the same governed IDs.",
        "Rebuild interactions and squared terms before re-scoring.",
        "Translate modeled hazard and survival movement into executive and technical evidence.",
    ]
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(199, 208, 213)
        paragraph.space_after = Pt(12)
    _ppt_add_footer(slide)

    _ppt_add_picture_slide(
        prs,
        "Persona Lifecycle Evidence",
        plots.persona_km,
        "Observed Kaplan-Meier curves; personas are descriptive segments, not CoxPH model inputs.",
    )
    _ppt_add_picture_slide(
        prs,
        "Model-Derived Risk Tiers",
        plots.risk_tier_km,
        "Observed Kaplan-Meier evidence across baseline partial-hazard tiers.",
    )
    _ppt_add_picture_slide(
        prs,
        "Same-Cohort Predicted Survival",
        plots.baseline_vs_scenario,
        "CoxPH-predicted survival for the identical baseline target cohort under controlled scenarios.",
    )

    # Slide 5 — Scenario comparison table and chart.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_background(slide)
    _ppt_add_title(
        slide,
        "Scenario Comparison",
        f"Modeled sensitivity at a {config.evaluation_horizon:g}-unit evaluation horizon",
    )
    _ppt_add_scenario_table(slide, scenarios.summary)
    slide.shapes.add_picture(
        str(plots.scenario_comparison),
        Inches(1.55),
        Inches(4.20),
        width=Inches(10.3),
    )
    _ppt_add_footer(slide)

    # Slide 6 — Executive interpretation and boundaries.
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ppt_set_background(slide)
    _ppt_add_title(slide, "Executive Interpretation", "Model evidence with explicit boundaries")
    narrative_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.8), Inches(2.4))
    frame = narrative_box.text_frame
    frame.word_wrap = True
    frame.text = narrative
    frame.paragraphs[0].font.name = "Arial"
    frame.paragraphs[0].font.size = Pt(16)
    frame.paragraphs[0].font.color.rgb = RGBColor(244, 245, 246)

    boundary_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.15), Inches(11.8), Inches(1.8))
    boundary_frame = boundary_box.text_frame
    boundaries = [
        "Modeled sensitivity — not causal treatment effect.",
        "Synthetic / public demonstration data — no PII.",
        "No booked revenue, realized ROI, or production forecast is claimed.",
        "Production use requires domain calibration, monitoring, and model-risk approval.",
    ]
    for index, boundary in enumerate(boundaries):
        paragraph = boundary_frame.paragraphs[0] if index == 0 else boundary_frame.add_paragraph()
        paragraph.text = boundary
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(217, 164, 65)
        paragraph.space_after = Pt(8)
    _ppt_add_footer(slide)

    prs.save(output_path)


# ==================================================================================================
# TECHNICAL PDF
# ==================================================================================================

def dataframe_for_display(df: pd.DataFrame, max_rows: int = 30) -> pd.DataFrame:
    display = df.head(max_rows).copy()
    for column in display.columns:
        if pd.api.types.is_float_dtype(display[column]):
            if any(token in column.lower() for token in ["rate", "reduction", "uplift", "survival"]):
                display[column] = display[column].map(
                    lambda value: "" if pd.isna(value) else f"{value:.2%}"
                )
            elif "p_value" in column.lower():
                display[column] = display[column].map(
                    lambda value: "" if pd.isna(value) else f"{value:.4g}"
                )
            else:
                display[column] = display[column].map(
                    lambda value: "" if pd.isna(value) else f"{value:,.4f}"
                )
    return display


def reportlab_table_from_df(
    df: pd.DataFrame,
    max_rows: int = 30,
    font_size: float = 6.5,
) -> Table:
    display = dataframe_for_display(df, max_rows=max_rows)
    data: list[list[Any]] = [
        [Paragraph(str(column), ParagraphStyle("Header", fontSize=font_size, textColor=colors.white))
         for column in display.columns]
    ]
    for _, row in display.iterrows():
        data.append(
            [
                Paragraph(
                    "" if pd.isna(value) else str(value),
                    ParagraphStyle("Cell", fontSize=font_size, leading=font_size + 1.2),
                )
                for value in row.tolist()
            ]
        )

    usable_width = landscape(letter)[0] - 0.7 * inch
    column_width = usable_width / max(1, len(display.columns))
    table = Table(data, repeatRows=1, colWidths=[column_width] * len(display.columns))
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#14566A")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#F4F6F7")),
                ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#87949C")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 3),
                ("RIGHTPADDING", (0, 0), (-1, -1), 3),
                ("TOPPADDING", (0, 0), (-1, -1), 3),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
            ]
        )
    )
    return table


def export_technical_pdf(
    output_path: Path,
    config: FrameworkConfig,
    validation: pd.DataFrame,
    model: CoxModelArtifacts,
    personas: PersonaArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
    plots: PlotArtifacts,
    narrative: str,
) -> None:
    document = SimpleDocTemplate(
        str(output_path),
        pagesize=landscape(letter),
        leftMargin=0.35 * inch,
        rightMargin=0.35 * inch,
        topMargin=0.35 * inch,
        bottomMargin=0.35 * inch,
    )
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="CenteredTitle",
            parent=styles["Title"],
            alignment=TA_CENTER,
            textColor=colors.HexColor("#14566A"),
        )
    )
    styles.add(
        ParagraphStyle(
            name="EvidenceHeading",
            parent=styles["Heading2"],
            textColor=colors.HexColor("#14566A"),
            spaceBefore=8,
            spaceAfter=6,
        )
    )
    story: list[Any] = []

    story.append(Paragraph("Technical Model Evidence", styles["CenteredTitle"]))
    story.append(Paragraph(config.project_name, styles["Heading2"]))
    story.append(Paragraph(narrative, styles["BodyText"]))
    story.append(Spacer(1, 10))

    story.append(Paragraph("1. Input Validation", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(validation, max_rows=40, font_size=7))

    story.append(Paragraph("2. Model Fit Metadata", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(model.fit_metadata, max_rows=10, font_size=7))

    story.append(Paragraph("3. CoxPH Model Summary", styles["EvidenceHeading"]))
    model_columns = [
        column
        for column in [
            "Feature",
            "Coefficient",
            "Hazard_Ratio",
            "Hazard_Ratio_CI_Lower",
            "Hazard_Ratio_CI_Upper",
            "P_Value",
            "Interpretation_Unit",
        ]
        if column in model.model_summary.columns
    ]
    story.append(reportlab_table_from_df(model.model_summary[model_columns], max_rows=40, font_size=6.5))

    story.append(PageBreak())
    story.append(Paragraph("4. Proportional-Hazards Test", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(model.ph_test_results, max_rows=40, font_size=7))

    story.append(Paragraph("5. Persona Profiles", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(personas.profiles, max_rows=20, font_size=6.5))
    story.append(Spacer(1, 8))
    story.append(RLImage(str(plots.persona_km), width=5.2 * inch, height=3.1 * inch))

    story.append(PageBreak())
    story.append(Paragraph("6. Risk-Tier Profiles", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(risk.risk_tier_profiles, max_rows=20, font_size=6.5))
    story.append(Spacer(1, 8))
    story.append(RLImage(str(plots.risk_tier_km), width=5.2 * inch, height=3.1 * inch))

    story.append(PageBreak())
    story.append(Paragraph("7. Scenario Results", styles["EvidenceHeading"]))
    scenario_columns = [
        "Scenario_Rank",
        "Scenario_Name",
        "Target_Cohort_Count",
        "Relative_Hazard_Reduction",
        "Baseline_Survival_At_Horizon",
        "Scenario_Survival_At_Horizon",
        "Survival_Probability_Uplift",
        "Modeled_Direction",
    ]
    story.append(reportlab_table_from_df(scenarios.summary[scenario_columns], max_rows=20, font_size=6.5))
    story.append(Spacer(1, 8))
    story.append(RLImage(str(plots.baseline_vs_scenario), width=5.2 * inch, height=3.1 * inch))

    story.append(PageBreak())
    story.append(Paragraph("8. Dependency Synchronization Audit", styles["EvidenceHeading"]))
    story.append(reportlab_table_from_df(scenarios.dependency_audit, max_rows=60, font_size=6.2))

    story.append(Paragraph("9. Interpretation Boundaries", styles["EvidenceHeading"]))
    boundaries = [
        "The analysis uses synthetic or public demonstration data.",
        "Scenario movement is modeled sensitivity, not causal treatment effect.",
        "Partial hazard is relative risk, not an event probability.",
        "Predicted survival is a fitted-model output and requires validation before production use.",
        "No realized revenue, ROI, or defended LTV is claimed.",
    ]
    for boundary in boundaries:
        story.append(Paragraph(f"• {boundary}", styles["BodyText"]))

    document.build(story)


# ==================================================================================================
# RUN ARCHIVE, ACCEPTANCE CHECKS, AND MANIFEST
# ==================================================================================================

def run_acceptance_checks(
    df: pd.DataFrame,
    config: FrameworkConfig,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
) -> pd.DataFrame:
    checks: list[dict[str, Any]] = []

    def add(name: str, passed: bool, evidence: str) -> None:
        checks.append(
            {
                "Check": name,
                "Status": "PASS" if passed else "FAIL",
                "Evidence": evidence,
            }
        )

    add(
        "Unique governed IDs",
        not df[config.id_col].duplicated().any(),
        f"{df[config.id_col].nunique():,} unique IDs across {len(df):,} rows.",
    )
    add(
        "Target cohort is non-empty",
        len(risk.target_cohort) > 0,
        f"{len(risk.target_cohort):,} records selected at quantile {config.risk_quantile:.2f}.",
    )

    expected_ids = set(risk.target_cohort[config.id_col].astype(str))
    consistent_ids = True
    if not scenarios.target_scores.empty:
        for _, group in scenarios.target_scores.groupby("Scenario_Name"):
            if set(group[config.id_col].astype(str)) != expected_ids:
                consistent_ids = False
                break
    add(
        "Same target cohort across scenarios",
        consistent_ids,
        "Every scenario score table contains the identical baseline target IDs.",
    )

    dependency_pass = (
        scenarios.dependency_audit.empty
        or scenarios.dependency_audit["Validation_Status"].isin(["PASS", "NOT_APPLICABLE"]).all()
    )
    add(
        "Engineered dependencies synchronized",
        bool(dependency_pass),
        "All interaction and squared-term post-rebuild checks passed.",
    )

    finite_metrics = True
    if not scenarios.summary.empty:
        metric_columns = [
            "Baseline_Mean_Partial_Hazard",
            "Scenario_Mean_Partial_Hazard",
            "Relative_Hazard_Reduction",
            "Baseline_Survival_At_Horizon",
            "Scenario_Survival_At_Horizon",
            "Survival_Probability_Uplift",
        ]
        finite_metrics = np.isfinite(
            scenarios.summary[metric_columns].to_numpy(dtype=float)
        ).all()
    add(
        "Scenario metrics are calculated and finite",
        bool(finite_metrics),
        "No hard-coded uncertainty, ROI, revenue, or defended-value metric is used.",
    )

    curve_bounds = (
        scenarios.predicted_survival_curves.drop(columns=["Time"])
        .apply(lambda column: column.between(0.0, 1.0).all())
        .all()
    )
    add(
        "Predicted survival remains within [0, 1]",
        bool(curve_bounds),
        "All baseline and scenario survival probabilities satisfy probability bounds.",
    )

    return pd.DataFrame(checks)


def archive_run_tables(
    run_directory: Path,
    config: FrameworkConfig,
    scenarios_config: Sequence[ScenarioDefinition],
    input_df: pd.DataFrame,
    validation: pd.DataFrame,
    personas: PersonaArtifacts,
    model: CoxModelArtifacts,
    risk: RiskArtifacts,
    scenarios: ScenarioArtifacts,
    narrative: str,
) -> None:
    write_json(run_directory / "framework_config.json", config)
    write_json(run_directory / "scenario_definitions.json", list(scenarios_config))

    validation.to_csv(run_directory / "input_validation.csv", index=False)
    if config.save_input_snapshot:
        input_df.to_csv(run_directory / "input_snapshot.csv", index=False)

    model.model_summary.to_csv(run_directory / "model_summary.csv", index=False)
    model.fit_metadata.to_csv(run_directory / "model_fit_metadata.csv", index=False)
    model.ph_test_results.to_csv(run_directory / "proportional_hazards_test.csv", index=False)
    model.fit_warnings.to_csv(run_directory / "model_fit_warnings.csv", index=False)

    # Preserve the fitted transformation parameters used by both analytical branches.
    persona_scaler_parameters = pd.DataFrame(
        {
            "Feature": list(config.segmentation_features),
            "Mean": personas.scaler.mean_,
            "Scale": personas.scaler.scale_,
        }
    )
    persona_scaler_parameters.to_csv(
        run_directory / "persona_scaler_parameters.csv",
        index=False,
    )
    standardized_centroids = pd.DataFrame(
        personas.kmeans.cluster_centers_,
        columns=list(config.segmentation_features),
    )
    standardized_centroids.insert(
        0,
        "Persona",
        [f"Persona Group {index + 1}" for index in range(len(standardized_centroids))],
    )
    standardized_centroids.to_csv(
        run_directory / "persona_centroids_standardized.csv",
        index=False,
    )
    raw_centroids = pd.DataFrame(
        personas.scaler.inverse_transform(personas.kmeans.cluster_centers_),
        columns=list(config.segmentation_features),
    )
    raw_centroids.insert(
        0,
        "Persona",
        [f"Persona Group {index + 1}" for index in range(len(raw_centroids))],
    )
    raw_centroids.to_csv(
        run_directory / "persona_centroids_raw_scale.csv",
        index=False,
    )

    if model.model_scaler is not None:
        model_scaler_parameters = pd.DataFrame(
            {
                "Feature": list(config.model_features),
                "Mean": model.model_scaler.mean_,
                "Scale": model.model_scaler.scale_,
            }
        )
    else:
        model_scaler_parameters = pd.DataFrame(
            {
                "Feature": list(config.model_features),
                "Mean": np.nan,
                "Scale": np.nan,
                "Note": "Model features were not standardized.",
            }
        )
    model_scaler_parameters.to_csv(
        run_directory / "model_scaler_parameters.csv",
        index=False,
    )

    model.model.baseline_survival_.reset_index().to_csv(
        run_directory / "cox_baseline_survival.csv",
        index=False,
    )
    model.model.baseline_cumulative_hazard_.reset_index().to_csv(
        run_directory / "cox_baseline_cumulative_hazard.csv",
        index=False,
    )

    personas.profiles.to_csv(run_directory / "persona_profiles.csv", index=False)
    risk.risk_tier_profiles.to_csv(run_directory / "risk_tier_profiles.csv", index=False)
    risk.target_cohort.to_csv(run_directory / "target_cohort.csv", index=False)

    scenarios.summary.to_csv(run_directory / "scenario_results.csv", index=False)
    scenarios.target_scores.to_csv(run_directory / "scenario_target_scores.csv", index=False)
    scenarios.dependency_audit.to_csv(run_directory / "dependency_audit.csv", index=False)
    scenarios.predicted_survival_curves.to_csv(
        run_directory / "predicted_survival_curves.csv",
        index=False,
    )

    (run_directory / "executive_narrative.txt").write_text(narrative, encoding="utf-8")


def build_manifest(run_directory: Path) -> dict[str, str]:
    return {
        path.relative_to(run_directory).as_posix(): str(path.resolve())
        for path in sorted(run_directory.rglob("*"))
        if path.is_file()
    }


# ==================================================================================================
# FRAMEWORK ORCHESTRATOR
# ==================================================================================================

def run_survival_strategy_framework(
    df: pd.DataFrame,
    config: FrameworkConfig,
    scenarios: Sequence[ScenarioDefinition],
) -> RunArtifacts:
    """Executes the full governed workflow and returns the artifact package."""
    validation_result = validate_survival_input(df, config, scenarios)
    validation_result.raise_for_errors()
    validation_frame = validation_result.to_frame()

    fingerprint = dataframe_fingerprint(df, config.id_col)
    run_id, run_directory = create_run_directory(config, fingerprint)

    personas = fit_personas(df, config)
    model = fit_cox_model(df, config)
    risk = score_and_stratify_population(df, personas, model, config)
    scenario_artifacts = simulate_scenarios(
        df,
        config,
        scenarios,
        model,
        risk,
    )

    narrative = build_executive_narrative(config, model, risk, scenario_artifacts)
    archive_run_tables(
        run_directory,
        config,
        scenarios,
        df,
        validation_frame,
        personas,
        model,
        risk,
        scenario_artifacts,
        narrative,
    )

    plots = generate_evidence_plots(
        run_directory,
        config,
        personas,
        risk,
        scenario_artifacts,
    )

    power_point = run_directory / "Survival_Strategy_Deck.pptx"
    technical_pdf = run_directory / "Technical_Model_Evidence.pdf"
    export_executive_powerpoint(
        power_point,
        config,
        model,
        risk,
        scenario_artifacts,
        plots,
        narrative,
    )
    export_technical_pdf(
        technical_pdf,
        config,
        validation_frame,
        model,
        personas,
        risk,
        scenario_artifacts,
        plots,
        narrative,
    )

    acceptance_checks = run_acceptance_checks(
        df,
        config,
        risk,
        scenario_artifacts,
    )
    acceptance_checks.to_csv(run_directory / "acceptance_checks.csv", index=False)

    required_files = [power_point, technical_pdf]
    files_exist = all(path.exists() and path.stat().st_size > 0 for path in required_files)
    file_check = pd.DataFrame(
        [
            {
                "Check": "Executive and technical reports created",
                "Status": "PASS" if files_exist else "FAIL",
                "Evidence": " | ".join(path.name for path in required_files),
            }
        ]
    )
    acceptance_checks = pd.concat(
        [acceptance_checks, file_check],
        ignore_index=True,
    )
    acceptance_checks.to_csv(run_directory / "acceptance_checks.csv", index=False)

    run_metadata = {
        "engine_release": ENGINE_RELEASE,
        "run_id": run_id,
        "run_timestamp_utc": utc_now(),
        "project_name": config.project_name,
        "input_fingerprint_sha256": fingerprint,
        "row_count": len(df),
        "event_count": int(df[config.event_col].sum()),
        "event_rate": float(df[config.event_col].mean()),
        "duration_col": config.duration_col,
        "event_col": config.event_col,
        "id_col": config.id_col,
        "model_features": config.model_features,
        "segmentation_features": config.segmentation_features,
        "risk_quantile": config.risk_quantile,
        "target_cohort_count": len(risk.target_cohort),
        "penalizer": config.penalizer,
        "l1_ratio": config.l1_ratio,
        "random_state": config.random_state,
        "scenario_names": [scenario.name for scenario in scenarios],
        "package_versions": package_versions(),
        "acceptance_status": (
            "PASS"
            if (acceptance_checks["Status"] == "PASS").all()
            else "REVIEW"
        ),
    }
    write_json(run_directory / "run_metadata.json", run_metadata)

    manifest = build_manifest(run_directory)
    write_json(run_directory / "artifact_manifest.json", manifest)

    return RunArtifacts(
        run_id=run_id,
        run_directory=run_directory,
        validation=validation_frame,
        personas=personas,
        model=model,
        risk=risk,
        scenarios=scenario_artifacts,
        plots=plots,
        executive_narrative=narrative,
        acceptance_checks=acceptance_checks,
        power_point=power_point,
        technical_pdf=technical_pdf,
        manifest=manifest,
    )


# ==================================================================================================
# ROSSI SMOKE TEST — SECONDARY TECHNICAL DEMONSTRATION
# ==================================================================================================

def build_rossi_smoke_test(
    output_root: Path,
) -> tuple[pd.DataFrame, FrameworkConfig, tuple[ScenarioDefinition, ...]]:
    """Builds a secondary public-data smoke test using the lifelines Rossi dataset.

    Important boundary:
    The two scenarios below are technical sensitivity checks used to prove framework portability.
    They are not presented as retention interventions, causal strategies, or criminal-justice policy.
    """
    data = load_rossi().reset_index(drop=True)
    data.insert(0, "subject_id", [f"ROSSI-{i:04d}" for i in range(1, len(data) + 1)])
    data["age_x_prio"] = data["age"] * data["prio"]

    config = FrameworkConfig(
        project_name="Survival Strategy Framework - Rossi Technical Smoke Test",
        duration_col="week",
        event_col="arrest",
        id_col="subject_id",
        model_features=(
            "fin",
            "age",
            "race",
            "wexp",
            "mar",
            "paro",
            "prio",
            "age_x_prio",
        ),
        segmentation_features=("age", "prio", "fin"),
        profile_features=("age", "prio", "fin", "wexp", "mar"),
        n_clusters=3,
        risk_quantile=0.75,
        penalizer=0.10,
        random_state=42,
        standardize_model_features=True,
        show_km_confidence_intervals=True,
        run_ph_test=True,
        evaluation_horizon=26.0,
        timeline_points=105,
        output_root=output_root,
        save_input_snapshot=True,
    )

    scenarios = (
        ScenarioDefinition(
            name="Financial Support Flag Enabled - Technical Sensitivity",
            description=(
                "Set the public Rossi financial-support indicator to one for the same high-risk "
                "cohort as a technical portability test."
            ),
            changes=(
                FeatureChange(
                    feature="fin",
                    operation="replace",
                    value=1.0,
                    lower_bound=0.0,
                    upper_bound=1.0,
                    round_digits=0,
                ),
            ),
        ),
        ScenarioDefinition(
            name="Prior Incident Count Lower - Technical Sensitivity",
            description=(
                "Reduce the prior-incident count by one, bounded at zero, solely to verify "
                "dependency-safe same-cohort scoring in a second dataset."
            ),
            changes=(
                FeatureChange(
                    feature="prio",
                    operation="add",
                    value=-1.0,
                    lower_bound=0.0,
                    round_digits=0,
                ),
            ),
        ),
    )
    return data, config, scenarios


# ==================================================================================================
# BUILT-IN SELF-TESTS
# ==================================================================================================

def run_internal_self_tests() -> pd.DataFrame:
    """Runs lightweight unit-style checks without requiring pytest."""
    results: list[dict[str, str]] = []

    def record(name: str, passed: bool, detail: str) -> None:
        results.append(
            {
                "Test": name,
                "Status": "PASS" if passed else "FAIL",
                "Detail": detail,
            }
        )
        if not passed:
            raise AssertionError(f"Self-test failed: {name} — {detail}")

    data_a = generate_synthetic_retention_data(n_records=1_000, seed=17)
    data_b = generate_synthetic_retention_data(n_records=1_000, seed=17)
    record(
        "Synthetic data reproducibility",
        data_a.equals(data_b),
        "Same seed generated identical data.",
    )

    test_frame = pd.DataFrame(
        {
            "a": [1.0, 2.0],
            "b": [3.0, 4.0],
            "a_x_b": [0.0, 0.0],
            "a_sq": [0.0, 0.0],
        }
    )
    synced, audit = sync_engineered_dependencies(
        test_frame,
        ["a", "b", "a_x_b", "a_sq"],
        "SELF_TEST",
    )
    record(
        "Interaction dependency synchronization",
        np.allclose(synced["a_x_b"], test_frame["a"] * test_frame["b"]),
        "Interaction term was rebuilt from source fields.",
    )
    record(
        "Squared-term dependency synchronization",
        np.allclose(synced["a_sq"], test_frame["a"] ** 2),
        "Squared term was rebuilt from its base field.",
    )
    record(
        "Dependency audit status",
        audit["Validation_Status"].eq("PASS").all(),
        "Post-rebuild audit checks passed.",
    )

    bad_data = data_a.copy()
    bad_data.loc[0, "attrition_event"] = 2
    config = build_default_retention_config(Path("outputs_self_test"))
    validation = validate_survival_input(
        bad_data,
        config,
        build_default_retention_scenarios(),
    )
    record(
        "Nonbinary event rejection",
        validation.has_errors,
        "Invalid event value was rejected before modeling.",
    )

    return pd.DataFrame(results)


# ==================================================================================================
# COMMAND-LINE ENTRY POINT
# ==================================================================================================

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Run the Survival Strategy Framework demonstration.",
    )
    parser.add_argument(
        "--demo",
        choices=("retention", "rossi"),
        default="retention",
        help="Primary synthetic retention demo or secondary Rossi smoke test.",
    )
    parser.add_argument(
        "--records",
        type=int,
        default=7_500,
        help="Synthetic retention row count; ignored for Rossi.",
    )
    parser.add_argument(
        "--seed",
        type=int,
        default=42,
        help="Deterministic synthetic data seed.",
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=Path("outputs"),
        help="Root directory for immutable run packages.",
    )
    parser.add_argument(
        "--self-test",
        action="store_true",
        help="Run built-in unit-style checks before the framework run.",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()

    try:
        if args.self_test:
            self_test_results = run_internal_self_tests()
            print("\nINTERNAL SELF-TESTS")
            print(self_test_results.to_string(index=False))

        if args.demo == "retention":
            data = generate_synthetic_retention_data(
                n_records=args.records,
                seed=args.seed,
            )
            config = build_default_retention_config(args.output_root)
            scenarios = build_default_retention_scenarios()
        else:
            data, config, scenarios = build_rossi_smoke_test(args.output_root)

        artifacts = run_survival_strategy_framework(
            df=data,
            config=config,
            scenarios=scenarios,
        )

        print("\n" + "=" * 100)
        print("SURVIVAL STRATEGY FRAMEWORK — PIPELINE COMPLETE")
        print("=" * 100)
        print(f"Run ID:             {artifacts.run_id}")
        print(f"Run directory:      {artifacts.run_directory.resolve()}")
        print(f"Executive PPTX:     {artifacts.power_point.resolve()}")
        print(f"Technical PDF:      {artifacts.technical_pdf.resolve()}")
        print(f"Target cohort:      {len(artifacts.risk.target_cohort):,}")
        print("\nScenario results:")
        if artifacts.scenarios.summary.empty:
            print("No configured scenarios.")
        else:
            columns = [
                "Scenario_Rank",
                "Scenario_Name",
                "Relative_Hazard_Reduction",
                "Survival_Probability_Uplift",
            ]
            print(artifacts.scenarios.summary[columns].to_string(index=False))
        print("\nAcceptance checks:")
        print(artifacts.acceptance_checks.to_string(index=False))
        print("\nBoundary: modeled sensitivity, not causal or production impact.")
        return 0

    except Exception as exc:
        print("\nPIPELINE FAILED", file=sys.stderr)
        print(str(exc), file=sys.stderr)
        traceback.print_exc()
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
